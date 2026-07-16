"""
为 Fig1（SCOPE五层架构）和 Fig4（进化路线）生成可编辑的 PPTX 版本
每个形状、文本框、箭头都是独立对象，可在 PowerPoint 中自由拖拽调整
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ============================================================
# Fig 1: SCOPE五层架构与五步法 — 可编辑 PPTX
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.33)   # 16:9 wide
prs.slide_height = Inches(7.5)
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Color scheme
C_LAYERS = [
    RGBColor(0xE3, 0xF2, 0xFD),  # blue
    RGBColor(0xE8, 0xF5, 0xE9),  # green
    RGBColor(0xFF, 0xF3, 0xE0),  # orange
    RGBColor(0xF3, 0xE5, 0xF5),  # purple
    RGBColor(0xFF, 0xEB, 0xEE),  # red
]
C_LAYER_STROKE = [
    RGBColor(0x15, 0x65, 0xC0),
    RGBColor(0x2E, 0x7D, 0x32),
    RGBColor(0xE6, 0x51, 0x00),
    RGBColor(0x6A, 0x1B, 0x9A),
    RGBColor(0xC6, 0x28, 0x28),
]
C_STEP_FILL = [
    RGBColor(0x15, 0x65, 0xC0),
    RGBColor(0x02, 0x77, 0xBD),
    RGBColor(0x00, 0x83, 0x8F),
    RGBColor(0x00, 0x69, 0x5C),
    RGBColor(0x2E, 0x7D, 0x32),
]

def add_rounded_rect(slide, left, top, width, height, fill_color, stroke_color=None, stroke_width=Pt(1.5)):
    """Add a rounded rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if stroke_color:
        shape.line.color.rgb = stroke_color
        shape.line.width = stroke_width
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.08
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=Pt(11), bold=False,
                color=RGBColor(0x26, 0x32, 0x38), alignment=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_arrow(slide, start_left, start_top, end_left, end_top, color=RGBColor(0x54, 0x6E, 0x7A), width=Pt(1.5)):
    """Add a connector arrow"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.line.end_style = 'arrow'
    return connector

# === Title ===
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.5),
            '图1  SCOPE方法论五层架构与五步法流程',
            font_size=Pt(22), bold=True, color=RGBColor(0x1A, 0x23, 0x7E))
add_textbox(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.35),
            'Figure 1  Five-Layer Architecture and Five-Step Workflow of the SCOPE Methodology',
            font_size=Pt(13), bold=False, color=RGBColor(0x54, 0x6E, 0x7A))

# === LEFT SIDE: 5 Layers ===
layer_data = [
    ('第一层: 业务场景参数',
     '6个改善参数 + 6个恶化参数\n物理矛盾建模\n"想提升什么 vs 担心什么恶化"'),
    ('第二层: AI业务发明原理库',
     '36个原理, 五大类\n标准化五段模板\n基础原理 / 禁忌 / 信度标记'),
    ('第三层: 矛盾矩阵',
     '6×6 = 36格, 100%填满率\n210案例交叉统计\n信度等级 + 禁忌组合'),
    ('第四层: AI业务进化路线',
     '7条进化路线 + 回滚分支\n三级跳: 近期 → 中期 → 远期\n时间维度的推进策略'),
    ('第五层: AI能力效应库',
     '14项 AI技术映射\n"发明原理 → 具体AI技术"\nLLM / CV / RL / 时序预测 ...'),
]

layer_left = Inches(0.5)
layer_top_start = Inches(1.3)
layer_w = Inches(5.8)
layer_h = Inches(1.12)
layer_gap = Inches(0.1)

for i, (title, desc) in enumerate(layer_data):
    y = layer_top_start + i * (layer_h + layer_gap)

    # Layer box
    box = add_rounded_rect(slide, layer_left, y, layer_w, layer_h,
                          C_LAYERS[i], C_LAYER_STROKE[i], Pt(1.5))

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, layer_left + Inches(0.15), y + Inches(0.35),
        Inches(0.42), Inches(0.42)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_LAYER_STROKE[i]
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Title text
    add_textbox(slide, layer_left + Inches(0.75), y + Inches(0.08),
                layer_w - Inches(0.9), Inches(0.35),
                title, font_size=Pt(13), bold=True, color=RGBColor(0x26, 0x32, 0x38),
                alignment=PP_ALIGN.LEFT)

    # Description text
    add_textbox(slide, layer_left + Inches(0.75), y + Inches(0.38),
                layer_w - Inches(0.9), Inches(0.68),
                desc, font_size=Pt(10), bold=False, color=RGBColor(0x37, 0x47, 0x4F),
                alignment=PP_ALIGN.LEFT)

# Arrows between layers
for i in range(4):
    y_from = layer_top_start + i * (layer_h + layer_gap) + layer_h
    y_to = layer_top_start + (i + 1) * (layer_h + layer_gap)
    mid_x = layer_left + layer_w // 2
    add_arrow(slide, mid_x, y_from, mid_x, y_to,
              RGBColor(0x54, 0x6E, 0x7A), Pt(2.0))

# === RIGHT SIDE: 5 Steps ===
step_data = [
    ('S  Strengthen', '明确要提升的核心业务目标\n"我最想提升的是什么？"'),
    ('C  Constrain', '识别AI引入后的恶化风险\n"引入AI后，我最怕什么出问题？"'),
    ('O  Originate', '查矛盾矩阵\n获得AI业务发明原理推荐'),
    ('P  Path', '评估当前进化位置\n规划近期→中期→远期三级跳路径'),
    ('E  Effect', '匹配具体AI技术\n形成完整技术选型方案'),
]

step_left = Inches(7.2)
step_top_start = Inches(1.3)
step_w = Inches(5.6)
step_h = Inches(1.12)
step_gap = Inches(0.1)

for i, (title, desc) in enumerate(step_data):
    y = step_top_start + i * (step_h + step_gap)

    # Step box (white with colored border)
    box = add_rounded_rect(slide, step_left, y, step_w, step_h,
                          RGBColor(0xFF, 0xFF, 0xFF), C_STEP_FILL[i], Pt(2.5))

    # Title bar (colored)
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, step_left, y, step_w, Inches(0.35)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = C_STEP_FILL[i]
    title_bar.line.fill.background()

    # Step title
    add_textbox(slide, step_left + Inches(0.2), y + Inches(0.02),
                step_w - Inches(0.4), Inches(0.32),
                title, font_size=Pt(14), bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)

    # Step description
    add_textbox(slide, step_left + Inches(0.3), y + Inches(0.42),
                step_w - Inches(0.6), Inches(0.65),
                desc, font_size=Pt(11), bold=False,
                color=RGBColor(0x37, 0x47, 0x4F), alignment=PP_ALIGN.CENTER)

# Arrows between steps
for i in range(4):
    y_from = step_top_start + i * (step_h + step_gap) + step_h
    y_to = step_top_start + (i + 1) * (step_h + step_gap)
    mid_x = step_left + step_w // 2
    add_arrow(slide, mid_x, y_from, mid_x, y_to,
              RGBColor(0x45, 0x5A, 0x64), Pt(2.0))

# === Dashed connectors between layers and steps ===
link_labels = ['参数定义', '原理搜索', '矩阵查询', '路径评估', '技术匹配']
for i in range(5):
    ly = layer_top_start + i * (layer_h + layer_gap) + layer_h // 2
    sy = step_top_start + i * (step_h + step_gap) + step_h // 2

    # Dashed connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        layer_left + layer_w, ly, step_left, sy
    )
    connector.line.color.rgb = RGBColor(0x90, 0xA4, 0xAE)
    connector.line.width = Pt(1.0)
    connector.line.dash_style = 2  # dash

    # Label
    mid_y = (ly + sy) // 2
    add_textbox(slide, layer_left + layer_w + Inches(0.05), mid_y - Inches(0.2),
                Inches(0.85), Inches(0.35),
                link_labels[i], font_size=Pt(9), bold=False,
                color=RGBColor(0x78, 0x90, 0x9C), alignment=PP_ALIGN.CENTER)

# === Bottom annotation ===
add_textbox(slide, Inches(1), Inches(7.0), Inches(11.3), Inches(0.35),
            '注: 左侧五层为方法论的知识库支撑，右侧五步为用户交互流程，虚线连接表示每步对应的知识层调用。',
            font_size=Pt(10), bold=False, color=RGBColor(0x78, 0x90, 0x9C))

# Save Fig1 PPTX
fig1_path = os.path.join(OUTPUT_DIR, 'fig1_scope_architecture.pptx')
prs.save(fig1_path)
print(f"Fig 1 PPTX saved: {fig1_path}")

# ============================================================
# Fig 4: 七条AI业务进化路线 — 可编辑 PPTX
# ============================================================
prs2 = Presentation()
prs2.slide_width = Inches(13.33)
prs2.slide_height = Inches(7.5)
slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])

# Title
add_textbox(slide2, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.45),
            '图4  SCOPE七条AI业务进化路线',
            font_size=Pt(22), bold=True, color=RGBColor(0x1A, 0x23, 0x7E))
add_textbox(slide2, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.35),
            'Figure 4  Seven AI Business Evolution Trajectories of SCOPE',
            font_size=Pt(13), bold=False, color=RGBColor(0x54, 0x6E, 0x7A))

ROUTES = [
    ('路线1: 洞察深度',
     ['报表\n(发生了什么)', '仪表盘\n(正在发生什么)', '预测性警报\n(将要发生什么)', '规范建议\n(该怎么办)', '自主行动'],
     5, (0x15, 0x65, 0xC0), '南方电网"大瓦特·天象": 从负荷报表到AI自主调度的全阶段贯通'),
    ('路线2: 决策自主度',
     ['AI给出\n建议', 'AI+人\n协同决策', 'AI增强人', 'AI替代人\n(特定任务)', 'AI自主执行\n人监督', 'AI完全\n自主'],
     6, (0x2E, 0x7D, 0x32), '多数处于"AI增强人"阶段 | 医疗最保守 | 电网/制造最激进'),
    ('路线3: 产品形态',
     ['流程\n自动化', '智能\n产品化', '生态\n平台化'],
     3, (0xE6, 0x51, 0x00), '杭州城市大脑: 1.0交通治堵 → 3.0"模型+智能体"'),
    ('路线4: 技术范式',
     ['基于规则', '基于模型', '基于生成', '基于\n自进化'],
     4, (0x6A, 0x1B, 0x9A), '主流: 模型→生成过渡期 | Darktrace/西门子: 自进化前沿'),
    ('路线5: 交互方式',
     ['菜单/\n表单', '自然\n语言', '多模态\n融合', '无感交互\n(环境感知)'],
     4, (0x00, 0x83, 0x8F), '滴滴: 菜单→NL | 美团无人车: 多模态 | Vantiq: 无感'),
    ('路线6: Agent架构',
     ['单Agent', '多Agent\n协同', 'Agent生态\n(自主协作)'],
     3, (0xC6, 0x28, 0x28), 'Figma/智联/绝味: 多Agent协同 | 2025-2026爆发趋势'),
    ('路线7: 价值分配',
     ['层级中介', '平台中介', 'AI使能的\n去中介化'],
     3, (0x37, 0x47, 0x4F), '遂溪仙品荔: 果农+AI→消费者 | 阿桂: 农户+AI→直销'),
]

route_top_start = Inches(1.0)
route_h = Inches(0.85)
route_gap = Inches(0.06)
label_left = Inches(0.3)
label_w = Inches(2.0)
box_start_left = Inches(2.5)
box_area_w = Inches(9.8)

for idx, (title, stages, n_stages, color_tuple, example) in enumerate(ROUTES):
    route_color = RGBColor(*color_tuple)
    y = route_top_start + idx * (route_h + route_gap)

    # Route label
    add_textbox(slide2, label_left, y + Inches(0.2), label_w, Inches(0.45),
                title, font_size=Pt(13), bold=True, color=route_color, alignment=PP_ALIGN.LEFT)

    # Stage boxes
    box_spacing = box_area_w / n_stages
    box_w = box_spacing * Inches(0.7) / Inches(1)

    for si, stage in enumerate(stages):
        bx = box_start_left + si * box_spacing
        bw = box_spacing * 0.72

        # Box
        box = add_rounded_rect(slide2, bx, y, bw, route_h,
                              RGBColor(0xFF, 0xFF, 0xFF), route_color, Pt(2.0))

        # Filled overlay (lighter for earlier stages)
        r, g, b = color_tuple
        t = si / max(n_stages - 1, 1) * 0.6  # 0 to 0.6 opacity
        af_r = int(r + (255 - r) * (1 - t))
        af_g = int(g + (255 - g) * (1 - t))
        af_b = int(b + (255 - b) * (1 - t))
        alpha_fill = RGBColor(af_r, af_g, af_b)
        fill_box = add_rounded_rect(slide2, bx, y, bw, route_h,
                                    alpha_fill, None, Pt(0))

        # Stage text
        is_final = (si == n_stages - 1)
        add_textbox(slide2, bx + Inches(0.05), y + Inches(0.08),
                    bw - Inches(0.1), route_h - Inches(0.16),
                    stage, font_size=Pt(10), bold=is_final,
                    color=RGBColor(0x26, 0x32, 0x38), alignment=PP_ALIGN.CENTER)

        # Arrow to next stage
        if si < n_stages - 1:
            arrow_start_x = bx + bw + Inches(0.02)
            arrow_end_x = bx + box_spacing - Inches(0.02)
            arrow_y = y + route_h // 2
            add_arrow(slide2, arrow_start_x, arrow_y, arrow_end_x, arrow_y,
                     route_color, Pt(1.8))

    # Example annotation
    add_textbox(slide2, box_start_left, y + route_h + Inches(0.02),
                box_area_w, Inches(0.25),
                f'代表案例: {example}', font_size=Pt(8), bold=False,
                color=RGBColor(0x78, 0x90, 0x9C), alignment=PP_ALIGN.LEFT)

# Save Fig4 PPTX
fig4_path = os.path.join(OUTPUT_DIR, 'fig4_evolution_routes.pptx')
prs2.save(fig4_path)
print(f"Fig 4 PPTX saved: {fig4_path}")

print("\nDone! Both PPTX files generated. Open in PowerPoint to edit each element individually.")
