"""
SCOPE 详细版矛盾矩阵 — 每格显示完整原理编号
同时生成 PNG/PDF + 可编辑 PPTX
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import FancyBboxPatch
import pandas as pd
import numpy as np
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Font setup ──
CJK = ['SimHei','Noto Sans CJK SC','Source Han Sans SC','Microsoft YaHei']
avail = [f.name for f in fm.fontManager.ttflist]
ZH = next((f for f in CJK if f in avail), 'sans-serif')
plt.rcParams.update({
    'font.family':'sans-serif','font.sans-serif':[ZH,'DejaVu Sans'],
    'axes.unicode_minus':False,'figure.dpi':150,'savefig.dpi':400,'savefig.bbox':'tight'
})

# ── Load data ──
df = pd.read_csv(os.path.join(OUTPUT_DIR, 'matrix_data.csv'))
row_order = ['S1_运营效率','S2_客户满意度','S3_创新能力','S4_风险控制','S5_收入增长','S6_决策速度']
col_order = ['C1_实施复杂度','C2_可解释性','C3_对人的依赖度','C4_成本','C5_数据隐私风险','C6_系统安全性']

pvt_n = df.pivot_table(index='Improvement',columns='Constraint',values='CaseCount',aggfunc='first').reindex(index=row_order,columns=col_order)
pvt_c = df.pivot_table(index='Improvement',columns='Constraint',values='Confidence',aggfunc='first').reindex(index=row_order,columns=col_order)
pvt_p = df.pivot_table(index='Improvement',columns='Constraint',values='Principles',aggfunc='first').reindex(index=row_order,columns=col_order)

CONF_COLORS = {'High':'#1B7837','Medium':'#E6C800','Low':'#E68A00','Minimal':'#D73027'}
CONF_SYMS  = {'High':'■','Medium':'□','Low':'△','Minimal':'▽'}
conf_map   = {'High':3,'Medium':2,'Low':1,'Minimal':0}

color_mat = pvt_c.map(lambda x: conf_map.get(x,-1)).values.astype(float)
cmap = matplotlib.colors.ListedColormap([CONF_COLORS['Minimal'],CONF_COLORS['Low'],CONF_COLORS['Medium'],CONF_COLORS['High']])

# ── Build cell text with principle numbers (multi-line, word-wrapped) ──
def format_cell(i, j):
    """Format cell text: confidence line + principle number lines"""
    n_val = int(pvt_n.iloc[i,j])
    conf = pvt_c.iloc[i,j]
    sym = CONF_SYMS.get(conf, '?')
    principles_str = str(pvt_p.iloc[i,j]) if pd.notna(pvt_p.iloc[i,j]) else ''

    lines = []
    # Line 1: confidence + N
    lines.append(f'{sym} N={n_val}')
    # Line 2+: principle numbers, wrapped at ~30 chars per line
    if principles_str:
        # Split principle numbers
        nums = [p.strip() for p in principles_str.split(',') if p.strip()]
        # Add basic principles [2,7] note
        full_nums = nums + ['+2','+7']
        # Build wrapped lines
        current_line = ''
        for num in full_nums:
            test = current_line + (', ' if current_line else '') + num
            if len(test) > 28 and current_line:
                lines.append(current_line)
                current_line = num
            else:
                current_line = test
        if current_line:
            lines.append(current_line)
    return '\n'.join(lines)

# ── FIGURE: Large canvas ──
# Need VERY large cells for principle numbers
# Cell size: ~1.6 inch wide, ~1.5 inch tall → figure ~11.5 x 11 inches
CELL_W = 1.65
CELL_H = 1.55
FIG_W = CELL_W * 6 + 3.0   # + margin for labels & colorbar
FIG_H = CELL_H * 6 + 2.8   # + margin for titles

fig, ax = plt.subplots(figsize=(FIG_W, FIG_H))
ax.set_xlim(0, 6)
ax.set_ylim(0, 6)
ax.axis('off')

# Draw colored cells
for i in range(6):
    for j in range(6):
        cval = color_mat[i][j]
        color = cmap(cval / 3.0)
        rect = FancyBboxPatch((j, 5-i), 1, 1, boxstyle='round,pad=0.04',
                              facecolor=color, edgecolor='white', linewidth=2.0)
        ax.add_patch(rect)

        # Cell text
        cell_text = format_cell(i, j)
        text_color = 'white' if cval == 3 else 'black'
        # Split into lines and position each
        text_lines = cell_text.split('\n')
        total_lines = len(text_lines)
        line_h = 0.85 / max(total_lines, 1)  # dist between lines
        start_y = 5 - i + 0.5 + (total_lines - 1) * line_h / 2

        for li, line in enumerate(text_lines):
            y_pos = start_y - li * line_h
            fs = 8.5 if li == 0 else 6.8  # first line bigger
            fw = 'bold' if li == 0 else 'normal'
            ax.text(j + 0.5, y_pos, line, ha='center', va='center',
                   fontsize=fs, fontweight=fw, color=text_color)

# Axis labels
short_rows = ['S1\n运营效率','S2\n客户满意度','S3\n创新能力','S4\n风险控制','S5\n收入增长','S6\n决策速度']
short_cols = ['C1\n实施复杂度','C2\n可解释性','C3\n对人的\n依赖度','C4\n成本','C5\n数据\n隐私风险','C6\n系统\n安全性']

for i in range(6):
    ax.text(-0.55, 5-i+0.5, short_rows[i], ha='center', va='center',
           fontsize=10, fontweight='bold', color='#263238')
for j in range(6):
    ax.text(j+0.5, 6.35, short_cols[j], ha='center', va='center',
           fontsize=9, fontweight='bold', color='#263238')

# Axis titles
ax.text(3, -0.7, '恶化的业务约束 (Constrain) →', ha='center', va='center',
       fontsize=13, fontweight='bold', color='#1A237E')
ax.text(-1.25, 3, '← 改善的业务目标 (Strengthen)', ha='center', va='center',
       fontsize=13, fontweight='bold', color='#1A237E', rotation=90)

# ── Colorbar (below) ──
cbar_ax = fig.add_axes([0.18, -0.03, 0.64, 0.025])
cb = matplotlib.colorbar.ColorbarBase(cbar_ax, cmap=cmap, orientation='horizontal',
                                       ticks=[0.375, 1.125, 1.875, 2.625])
cb.set_ticklabels(['最低 (N=1-2)', '低 (N=3-5)', '中 (N=6-14)', '高 (N≥15)'])
cb.ax.tick_params(labelsize=8)

# Legend for symbols
sym_text = '符号: ■ 高信度  □ 中信度  △ 低信度  ▽ 最低信度  |  基础原理 [+2 认知自动化] [+7 预测与推演] 每格默认'
fig.text(0.5, 0.005, sym_text, ha='center', va='center', fontsize=7.5, color='#546E7A')

# Titles
fig.suptitle('SCOPE矛盾矩阵（详细版）— 每格标注关联的AI业务发明原理编号',
             fontsize=16, fontweight='bold', color='#1A237E', y=0.99)
fig.text(0.5, 0.955, 'SCOPE Contradiction Matrix (Detailed) — Principle Numbers Per Cell',
        ha='center', va='center', fontsize=10, style='italic', color='#546E7A')

# Save PNG + PDF
fig.savefig(os.path.join(OUTPUT_DIR, 'fig2_detailed_matrix.png'), dpi=400, facecolor='white')
fig.savefig(os.path.join(OUTPUT_DIR, 'fig2_detailed_matrix.pdf'), facecolor='white')
print("fig2_detailed_matrix.png + .pdf saved")
plt.close()

# ═══════════════════════════════════════════════════════════
# PPTX 可编辑版本
# ═══════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(13)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Color map for PPTX
PPTX_CONF = {
    'High': PptRGB(0x1B,0x78,0x37),
    'Medium': PptRGB(0xE6,0xC8,0x00),
    'Low': PptRGB(0xE6,0x8A,0x00),
    'Minimal': PptRGB(0xD7,0x30,0x27),
}

# Layout
CELL_W_PPTX = Inches(1.8)
CELL_H_PPTX = Inches(1.65)
MATRIX_LEFT = Inches(1.9)
MATRIX_TOP = Inches(1.6)
CELL_GAP = Inches(0.08)

def add_textbox_pptx(slide, left, top, width, height, text, font_size=Pt(8),
                     bold=False, color=PptRGB(0,0,0), alignment=PP_ALIGN.CENTER,
                     font_name='Microsoft YaHei'):
    """Add editable text box to slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

# Title
add_textbox_pptx(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.6),
    'SCOPE矛盾矩阵（详细版）— 每格标注关联的AI业务发明原理编号',
    Pt(22), True, PptRGB(0x1A,0x23,0x7E))
add_textbox_pptx(slide, Inches(1), Inches(0.65), Inches(14), Inches(0.35),
    'SCOPE Contradiction Matrix (Detailed) — Principle Numbers Per Cell  |  +2=认知自动化  +7=预测与推演（基础原理，每格默认）',
    Pt(11), False, PptRGB(0x54,0x6E,0x7A))

# Draw cells
for i in range(6):
    for j in range(6):
        x = MATRIX_LEFT + j * (CELL_W_PPTX + CELL_GAP)
        y = MATRIX_TOP + i * (CELL_H_PPTX + CELL_GAP)
        conf = pvt_c.iloc[i,j]
        cell_color = PPTX_CONF.get(conf, PptRGB(0xCC,0xCC,0xCC))

        # Cell rect
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, CELL_W_PPTX, CELL_H_PPTX)
        shape.fill.solid()
        shape.fill.fore_color.rgb = cell_color
        shape.line.color.rgb = PptRGB(0xFF,0xFF,0xFF)
        shape.line.width = Pt(2)
        shape.adjustments[0] = 0.05

        # Cell text
        cell_text = format_cell(i, j)
        text_color = PptRGB(0xFF,0xFF,0xFF) if conf == 'High' else PptRGB(0,0,0)
        add_textbox_pptx(slide, x + Inches(0.06), y + Inches(0.08),
                        CELL_W_PPTX - Inches(0.12), CELL_H_PPTX - Inches(0.16),
                        cell_text, Pt(7.5), False, text_color)

# Row labels
for i in range(6):
    y = MATRIX_TOP + i * (CELL_H_PPTX + CELL_GAP) + CELL_H_PPTX / 2
    add_textbox_pptx(slide, Inches(0.1), y - Inches(0.3), Inches(1.6), Inches(0.6),
        short_rows[i].replace('\n',' '), Pt(11), True, PptRGB(0x26,0x32,0x38))

# Column labels
for j in range(6):
    x = MATRIX_LEFT + j * (CELL_W_PPTX + CELL_GAP) + CELL_W_PPTX / 2
    add_textbox_pptx(slide, x - Inches(0.7), MATRIX_TOP - Inches(0.55), Inches(1.4), Inches(0.5),
        short_cols[j].replace('\n',' '), Pt(10), True, PptRGB(0x26,0x32,0x38))

# Axis titles
add_textbox_pptx(slide, Inches(4), Inches(0.95), Inches(8), Inches(0.4),
    '恶化的业务约束 (Constrain) →', Pt(14), True, PptRGB(0x1A,0x23,0x7E))
add_textbox_pptx(slide, Inches(0.05), Inches(4.5), Inches(1.5), Inches(1.5),
    '← 改善的业务目标\n   (Strengthen)', Pt(14), True, PptRGB(0x1A,0x23,0x7E))

# Legend area at bottom
legend_y = MATRIX_TOP + 6 * (CELL_H_PPTX + CELL_GAP) + Inches(0.3)
legend_items = [
    ('■ 高信度 (N≥15)', PptRGB(0x1B,0x78,0x37)),
    ('□ 中信度 (N=6-14)', PptRGB(0xE6,0xC8,0x00)),
    ('△ 低信度 (N=3-5)', PptRGB(0xE6,0x8A,0x00)),
    ('▽ 最低信度 (N=1-2)', PptRGB(0xD7,0x30,0x27)),
]
for k, (label, color) in enumerate(legend_items):
    lx = Inches(2 + k * 3.2)
    # Color swatch
    swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, legend_y, Inches(0.25), Inches(0.2))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = color
    swatch.line.fill.background()
    add_textbox_pptx(slide, lx + Inches(0.3), legend_y - Inches(0.02), Inches(2.5), Inches(0.25),
                    label, Pt(10), False, PptRGB(0x37,0x47,0x4F), PP_ALIGN.LEFT)

pptx_path = os.path.join(OUTPUT_DIR, 'fig2_detailed_matrix.pptx')
prs.save(pptx_path)
print(f"fig2_detailed_matrix.pptx saved")
print("Done!")
